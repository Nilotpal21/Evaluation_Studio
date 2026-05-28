"""
Comprehensive Format Tests for Docling Service

Tests ALL supported MIME types defined in document routing strategy:
- PDF (application/pdf)
- DOCX (application/vnd.openxmlformats-officedocument.wordprocessingml.document)
- DOC (application/msword)
- PPTX (application/vnd.openxmlformats-officedocument.presentationml.presentation)
- PPT (application/vnd.ms-powerpoint)
- HTML (text/html)
- Images: PNG, JPEG, TIFF, BMP, WEBP

Usage:
    pytest test_all_formats.py -v
    pytest test_all_formats.py -v -k "test_pdf"  # Run specific format
    pytest test_all_formats.py -v --tb=short     # Short traceback
"""

import pytest
import requests
import json
import io
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
import tempfile
from typing import Dict, Any


# ─── Configuration ───────────────────────────────────────────────────────────

SERVICE_URL = "http://localhost:8080"
TIMEOUT = 300  # 5 minutes


# ─── Test Document Generators ────────────────────────────────────────────────

class TestDocumentGenerator:
    """Generate test documents for each format"""

    @staticmethod
    def create_sample_text() -> str:
        """Sample text with structure for testing"""
        return """# Introduction to Document Processing

This is a sample document to test Docling extraction capabilities.

## Features

Docling supports multiple document formats:
- PDF documents
- Microsoft Office files (DOCX, PPTX)
- HTML content
- Images with OCR

## Table Example

| Format | Support | Quality |
|--------|---------|---------|
| PDF    | ✓       | High    |
| DOCX   | ✓       | High    |
| HTML   | ✓       | Medium  |

## Conclusion

This document tests basic text extraction, layout preservation, and table parsing.
"""

    @staticmethod
    def create_pdf_document() -> bytes:
        """Create a simple PDF for testing"""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle, Spacer
            from reportlab.lib import colors

            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []

            # Title
            story.append(Paragraph("Test PDF Document", styles['Title']))
            story.append(Spacer(1, 12))

            # Heading
            story.append(Paragraph("Introduction", styles['Heading1']))
            story.append(Paragraph(
                "This is a test PDF document created with ReportLab to validate Docling extraction.",
                styles['Normal']
            ))
            story.append(Spacer(1, 12))

            # Table
            story.append(Paragraph("Test Table", styles['Heading2']))
            data = [
                ['Format', 'Extension', 'Support'],
                ['PDF', '.pdf', 'Full'],
                ['DOCX', '.docx', 'Full'],
                ['HTML', '.html', 'Full'],
            ]
            table = Table(data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ]))
            story.append(table)

            doc.build(story)
            return buffer.getvalue()

        except ImportError:
            pytest.skip("reportlab not installed - pip install reportlab")

    @staticmethod
    def create_docx_document() -> bytes:
        """Create a simple DOCX for testing"""
        try:
            from docx import Document
            from docx.shared import Inches

            doc = Document()

            # Title
            doc.add_heading('Test DOCX Document', 0)

            # Paragraphs
            doc.add_heading('Introduction', level=1)
            doc.add_paragraph('This is a test DOCX document created with python-docx.')

            # Table
            doc.add_heading('Test Table', level=2)
            table = doc.add_table(rows=4, cols=3)
            table.style = 'Light Grid Accent 1'

            # Header row
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = 'Format'
            hdr_cells[1].text = 'Extension'
            hdr_cells[2].text = 'Support'

            # Data rows
            data_rows = [
                ('PDF', '.pdf', 'Full'),
                ('DOCX', '.docx', 'Full'),
                ('HTML', '.html', 'Full'),
            ]
            for i, (fmt, ext, support) in enumerate(data_rows, start=1):
                row_cells = table.rows[i].cells
                row_cells[0].text = fmt
                row_cells[1].text = ext
                row_cells[2].text = support

            # Save to bytes
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.read()

        except ImportError:
            pytest.skip("python-docx not installed - pip install python-docx")

    @staticmethod
    def create_pptx_document() -> bytes:
        """Create a simple PPTX for testing"""
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()

            # Slide 1: Title
            slide_layout = prs.slide_layouts[0]  # Title slide
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Test PPTX Document"
            subtitle.text = "Created with python-pptx"

            # Slide 2: Content with bullet points
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Supported Formats"

            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = "Docling supports:"

            p = tf.add_paragraph()
            p.text = "PDF documents"
            p.level = 1

            p = tf.add_paragraph()
            p.text = "Microsoft Office files"
            p.level = 1

            p = tf.add_paragraph()
            p.text = "HTML content"
            p.level = 1

            # Save to bytes
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return buffer.read()

        except ImportError:
            pytest.skip("python-pptx not installed - pip install python-pptx")

    @staticmethod
    def create_html_document() -> bytes:
        """Create a simple HTML document for testing"""
        html = """<!DOCTYPE html>
<html>
<head>
    <title>Test HTML Document</title>
</head>
<body>
    <h1>Introduction to Document Processing</h1>
    <p>This is a test HTML document to validate Docling extraction.</p>

    <h2>Features</h2>
    <ul>
        <li>PDF documents</li>
        <li>Microsoft Office files</li>
        <li>HTML content</li>
        <li>Images with OCR</li>
    </ul>

    <h2>Table Example</h2>
    <table border="1">
        <tr>
            <th>Format</th>
            <th>Extension</th>
            <th>Support</th>
        </tr>
        <tr>
            <td>PDF</td>
            <td>.pdf</td>
            <td>Full</td>
        </tr>
        <tr>
            <td>DOCX</td>
            <td>.docx</td>
            <td>Full</td>
        </tr>
        <tr>
            <td>HTML</td>
            <td>.html</td>
            <td>Full</td>
        </tr>
    </table>

    <h2>Conclusion</h2>
    <p>This document tests HTML parsing and structure extraction.</p>
</body>
</html>"""
        return html.encode('utf-8')

    @staticmethod
    def create_image_with_text(format: str = 'PNG') -> bytes:
        """Create an image with text for OCR testing"""
        # Create image
        img = Image.new('RGB', (800, 600), color='white')
        draw = ImageDraw.Draw(img)

        # Draw text
        try:
            # Try to use a nice font
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 40)
        except:
            # Fall back to default
            font = ImageFont.load_default()

        text = "Test Image with Text\n\nThis is a sample image\nfor OCR testing.\n\nDocling should extract\nthis text."
        draw.text((50, 50), text, fill='black', font=font)

        # Draw a simple table-like structure
        draw.rectangle([50, 350, 750, 550], outline='black', width=2)
        draw.line([50, 400, 750, 400], fill='black', width=2)
        draw.line([400, 350, 400, 550], fill='black', width=2)

        # Save to bytes
        buffer = io.BytesIO()
        img.save(buffer, format=format)
        buffer.seek(0)
        return buffer.read()


# ─── Fixtures ────────────────────────────────────────────────────────────────

@pytest.fixture(scope="session")
def service_url():
    """Docling service URL"""
    return SERVICE_URL


@pytest.fixture(scope="session")
def service_health(service_url):
    """Check if service is healthy before running tests"""
    try:
        response = requests.get(f"{service_url}/health", timeout=5)
        if response.status_code != 200:
            pytest.skip(f"Service not available (status {response.status_code})")

        health = response.json()
        if not health.get('docling_available'):
            pytest.skip("Docling library not available in service")

        return health
    except requests.exceptions.RequestException as e:
        pytest.skip(f"Cannot connect to service: {e}. Start with: docker-compose up docling-service")


@pytest.fixture
def generator():
    """Test document generator"""
    return TestDocumentGenerator()


# ─── Helper Functions ────────────────────────────────────────────────────────

def extract_document(service_url: str, file_bytes: bytes, filename: str, mime_type: str) -> Dict[str, Any]:
    """Helper to extract a document"""
    files = {'file': (filename, io.BytesIO(file_bytes), mime_type)}
    options = json.dumps({
        'extractImages': True,
        'extractTables': True,
        'preserveLayout': True,
        'renderScreenshots': False,  # Disable screenshots for speed
        'ocrEnabled': True,
    })
    data = {'options': options}

    response = requests.post(
        f"{service_url}/extract",
        files=files,
        data=data,
        timeout=TIMEOUT,
    )

    assert response.status_code == 200, f"Extraction failed: {response.text}"
    return response.json()


# ─── PDF Tests ───────────────────────────────────────────────────────────────

def test_pdf_extraction(service_url, service_health, generator):
    """Test PDF extraction (application/pdf)"""
    pdf_bytes = generator.create_pdf_document()

    result = extract_document(service_url, pdf_bytes, 'test.pdf', 'application/pdf')

    # Validate structure
    assert 'pages' in result
    assert 'metadata' in result
    assert len(result['pages']) >= 1

    # Validate text extraction
    page = result['pages'][0]
    assert len(page['text']) > 0
    assert 'Test PDF Document' in page['text'] or 'test pdf document' in page['text'].lower()

    # Validate table extraction
    metadata = result['metadata']
    assert metadata['totalTables'] >= 1


# ─── DOCX Tests ──────────────────────────────────────────────────────────────

def test_docx_extraction(service_url, service_health, generator):
    """Test DOCX extraction (application/vnd.openxmlformats-officedocument.wordprocessingml.document)"""
    docx_bytes = generator.create_docx_document()

    result = extract_document(
        service_url,
        docx_bytes,
        'test.docx',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    )

    # Validate structure
    assert 'pages' in result
    assert 'metadata' in result
    assert len(result['pages']) >= 1

    # Validate text extraction
    page = result['pages'][0]
    assert len(page['text']) > 0
    assert 'Test DOCX Document' in page['text'] or 'test docx document' in page['text'].lower()

    # Validate table extraction
    metadata = result['metadata']
    assert metadata['totalTables'] >= 1


def test_doc_mime_type(service_url, service_health, generator):
    """Test DOC MIME type handling (application/msword)"""
    # Note: We create DOCX but test with .doc MIME type
    # Docling should handle it gracefully
    docx_bytes = generator.create_docx_document()

    result = extract_document(
        service_url,
        docx_bytes,
        'test.doc',
        'application/msword'
    )

    # Should still extract successfully (Docling is format-agnostic)
    assert 'pages' in result
    assert len(result['pages']) >= 1


# ─── PPTX Tests ──────────────────────────────────────────────────────────────

def test_pptx_extraction(service_url, service_health, generator):
    """Test PPTX extraction (application/vnd.openxmlformats-officedocument.presentationml.presentation)"""
    pptx_bytes = generator.create_pptx_document()

    result = extract_document(
        service_url,
        pptx_bytes,
        'test.pptx',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

    # Validate structure
    assert 'pages' in result
    assert 'metadata' in result

    # PPTX should have 2 slides (pages)
    assert len(result['pages']) >= 2

    # Validate text extraction from first slide
    first_slide = result['pages'][0]
    assert len(first_slide['text']) > 0
    assert 'Test PPTX Document' in first_slide['text'] or 'test pptx' in first_slide['text'].lower()


def test_ppt_mime_type(service_url, service_health, generator):
    """Test PPT MIME type handling (application/vnd.ms-powerpoint)"""
    pptx_bytes = generator.create_pptx_document()

    result = extract_document(
        service_url,
        pptx_bytes,
        'test.ppt',
        'application/vnd.ms-powerpoint'
    )

    # Should still extract successfully
    assert 'pages' in result
    assert len(result['pages']) >= 1


# ─── HTML Tests ──────────────────────────────────────────────────────────────

def test_html_extraction(service_url, service_health, generator):
    """Test HTML extraction (text/html)"""
    html_bytes = generator.create_html_document()

    result = extract_document(service_url, html_bytes, 'test.html', 'text/html')

    # Validate structure
    assert 'pages' in result
    assert 'metadata' in result
    assert len(result['pages']) >= 1

    # Validate text extraction
    page = result['pages'][0]
    assert len(page['text']) > 0
    assert 'Introduction to Document Processing' in page['text'] or 'introduction' in page['text'].lower()

    # HTML should extract table
    metadata = result['metadata']
    assert metadata['totalTables'] >= 1


# ─── Image Tests ─────────────────────────────────────────────────────────────

def test_png_extraction(service_url, service_health, generator):
    """Test PNG image extraction with OCR (image/png)"""
    png_bytes = generator.create_image_with_text('PNG')

    result = extract_document(service_url, png_bytes, 'test.png', 'image/png')

    # Validate structure
    assert 'pages' in result
    assert 'metadata' in result
    assert len(result['pages']) >= 1

    # Validate OCR extraction
    page = result['pages'][0]
    assert len(page['text']) > 0
    # OCR should extract some text (may not be perfect)
    assert 'test' in page['text'].lower() or 'image' in page['text'].lower()


def test_jpeg_extraction(service_url, service_health, generator):
    """Test JPEG image extraction with OCR (image/jpeg)"""
    jpeg_bytes = generator.create_image_with_text('JPEG')

    result = extract_document(service_url, jpeg_bytes, 'test.jpg', 'image/jpeg')

    # Validate structure
    assert 'pages' in result
    assert len(result['pages']) >= 1
    assert len(result['pages'][0]['text']) > 0


def test_jpg_mime_type(service_url, service_health, generator):
    """Test alternative JPEG MIME type (image/jpg)"""
    jpeg_bytes = generator.create_image_with_text('JPEG')

    result = extract_document(service_url, jpeg_bytes, 'test.jpg', 'image/jpg')

    assert 'pages' in result
    assert len(result['pages']) >= 1


def test_tiff_extraction(service_url, service_health, generator):
    """Test TIFF image extraction with OCR (image/tiff)"""
    tiff_bytes = generator.create_image_with_text('TIFF')

    result = extract_document(service_url, tiff_bytes, 'test.tiff', 'image/tiff')

    # Validate structure
    assert 'pages' in result
    assert len(result['pages']) >= 1


def test_bmp_extraction(service_url, service_health, generator):
    """Test BMP image extraction with OCR (image/bmp)"""
    bmp_bytes = generator.create_image_with_text('BMP')

    result = extract_document(service_url, bmp_bytes, 'test.bmp', 'image/bmp')

    # Validate structure
    assert 'pages' in result
    assert len(result['pages']) >= 1


def test_webp_extraction(service_url, service_health, generator):
    """Test WebP image extraction with OCR (image/webp)"""
    webp_bytes = generator.create_image_with_text('WEBP')

    result = extract_document(service_url, webp_bytes, 'test.webp', 'image/webp')

    # Validate structure
    assert 'pages' in result
    assert len(result['pages']) >= 1


# ─── Comprehensive Format Coverage Test ──────────────────────────────────────

def test_all_supported_formats(service_url, service_health, generator):
    """Test that all formats defined in routing strategy are supported"""

    supported_formats = [
        ('PDF', 'application/pdf', 'test.pdf', generator.create_pdf_document),
        ('DOCX', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'test.docx', generator.create_docx_document),
        ('DOC', 'application/msword', 'test.doc', generator.create_docx_document),
        ('PPTX', 'application/vnd.openxmlformats-officedocument.presentationml.presentation', 'test.pptx', generator.create_pptx_document),
        ('PPT', 'application/vnd.ms-powerpoint', 'test.ppt', generator.create_pptx_document),
        ('HTML', 'text/html', 'test.html', generator.create_html_document),
        ('PNG', 'image/png', 'test.png', lambda: generator.create_image_with_text('PNG')),
        ('JPEG', 'image/jpeg', 'test.jpg', lambda: generator.create_image_with_text('JPEG')),
        ('JPG', 'image/jpg', 'test.jpg', lambda: generator.create_image_with_text('JPEG')),
        ('TIFF', 'image/tiff', 'test.tiff', lambda: generator.create_image_with_text('TIFF')),
        ('BMP', 'image/bmp', 'test.bmp', lambda: generator.create_image_with_text('BMP')),
        ('WEBP', 'image/webp', 'test.webp', lambda: generator.create_image_with_text('WEBP')),
    ]

    results = {}

    for format_name, mime_type, filename, generator_func in supported_formats:
        try:
            file_bytes = generator_func()
            result = extract_document(service_url, file_bytes, filename, mime_type)

            # Basic validation
            assert 'pages' in result, f"{format_name}: Missing 'pages' field"
            assert len(result['pages']) >= 1, f"{format_name}: No pages extracted"
            assert len(result['pages'][0]['text']) > 0, f"{format_name}: No text extracted"

            results[format_name] = '✅ PASS'

        except Exception as e:
            results[format_name] = f'❌ FAIL: {str(e)}'

    # Print summary
    print("\n" + "=" * 80)
    print("FORMAT SUPPORT SUMMARY")
    print("=" * 80)
    for format_name, status in results.items():
        print(f"{format_name:10} - {status}")
    print("=" * 80)

    # Assert all passed
    failures = [f"{name}: {status}" for name, status in results.items() if '❌' in status]
    assert len(failures) == 0, f"Some formats failed:\n" + "\n".join(failures)


# ─── Performance Benchmarks ──────────────────────────────────────────────────

@pytest.mark.benchmark
def test_extraction_performance(service_url, service_health, generator, benchmark):
    """Benchmark extraction performance"""
    pdf_bytes = generator.create_pdf_document()

    def extract():
        return extract_document(service_url, pdf_bytes, 'test.pdf', 'application/pdf')

    result = benchmark(extract)
    assert 'pages' in result


if __name__ == "__main__":
    pytest.main([__file__, "-v", "--tb=short"])
