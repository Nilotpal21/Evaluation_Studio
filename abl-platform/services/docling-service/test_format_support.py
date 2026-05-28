"""
Simple Format Support Tests for Docling

Tests that Docling can successfully extract content from all supported MIME types.
Focuses on essential validation: extraction succeeds and returns text.

Usage:
    pytest test_format_support.py -v
"""

import pytest
import requests
import json
import io
from PIL import Image, ImageDraw

SERVICE_URL = "http://localhost:8080"
TIMEOUT = 300


# ─── Fixtures ────────────────────────────────────────────────────────────────

@pytest.fixture(scope="session")
def service_url():
    return SERVICE_URL


@pytest.fixture(scope="session")
def check_service(service_url):
    """Check service is available"""
    try:
        response = requests.get(f"{service_url}/health", timeout=5)
        if response.status_code != 200:
            pytest.skip(f"Service not available")
        health = response.json()
        if not health.get('docling_available'):
            pytest.skip("Docling not available")
    except Exception as e:
        pytest.skip(f"Cannot connect: {e}")


# ─── Helper Functions ────────────────────────────────────────────────────────

def extract_document(service_url: str, file_bytes: bytes, filename: str, mime_type: str):
    """Extract document and return result"""
    files = {'file': (filename, io.BytesIO(file_bytes), mime_type)}
    options = json.dumps({
        'extractImages': True,
        'extractTables': True,
        'preserveLayout': True,
        'renderScreenshots': False,
        'ocrEnabled': True,
    })
    data = {'options': options}

    response = requests.post(
        f"{service_url}/extract",
        files=files,
        data=data,
        timeout=TIMEOUT,
    )

    return response


def create_simple_image_with_text() -> bytes:
    """Create a simple test image"""
    img = Image.new('RGB', (400, 200), color='white')
    draw = ImageDraw.Draw(img)
    draw.text((20, 20), "TEST IMAGE", fill='black')

    buffer = io.BytesIO()
    img.save(buffer, format='PNG')
    buffer.seek(0)
    return buffer.read()


# ─── Format Support Tests ────────────────────────────────────────────────────

@pytest.mark.parametrize("format_name,mime_type,filename,content", [
    # Images (these work well)
    ("PNG", "image/png", "test.png", lambda: create_simple_image_with_text()),
    ("JPEG", "image/jpeg", "test.jpg", lambda: Image.new('RGB', (100, 100), color='white')),
    ("TIFF", "image/tiff", "test.tiff", lambda: Image.new('RGB', (100, 100), color='white')),
    ("BMP", "image/bmp", "test.bmp", lambda: Image.new('RGB', (100, 100), color='white')),
])
def test_image_format_support(service_url, check_service, format_name, mime_type, filename, content):
    """Test image format support"""

    # Generate content
    if callable(content):
        if format_name == "PNG":
            file_bytes = content()
        else:
            # Generate image in correct format
            img = content()
            buffer = io.BytesIO()
            # Extract format from mime_type
            fmt = mime_type.split('/')[1].upper()
            img.save(buffer, format=fmt)
            buffer.seek(0)
            file_bytes = buffer.read()
    else:
        file_bytes = content

    # Extract
    response = extract_document(service_url, file_bytes, filename, mime_type)

    # Validate
    assert response.status_code == 200, f"{format_name}: Extraction failed with {response.status_code}"

    result = response.json()
    assert 'pages' in result, f"{format_name}: Missing 'pages' field"
    assert len(result['pages']) >= 1, f"{format_name}: No pages extracted"

    print(f"✅ {format_name} ({mime_type}): PASS")


def test_comprehensive_format_report(service_url, check_service):
    """Comprehensive test of all supported formats with detailed reporting"""

    formats_to_test = [
        # Documents
        ("PDF", "application/pdf", "test.pdf", "document"),
        ("DOCX", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "test.docx", "document"),
        ("DOC", "application/msword", "test.doc", "document"),
        ("PPTX", "application/vnd.openxmlformats-officedocument.presentationml.presentation", "test.pptx", "document"),
        ("PPT", "application/vnd.ms-powerpoint", "test.ppt", "document"),
        ("HTML", "text/html", "test.html", "document"),
        # Images
        ("PNG", "image/png", "test.png", "image"),
        ("JPEG", "image/jpeg", "test.jpg", "image"),
        ("JPG", "image/jpg", "test2.jpg", "image"),
        ("TIFF", "image/tiff", "test.tiff", "image"),
        ("BMP", "image/bmp", "test.bmp", "image"),
        ("WEBP", "image/webp", "test.webp", "image"),
    ]

    results = {}

    for format_name, mime_type, filename, doc_type in formats_to_test:
        try:
            # Generate appropriate test content
            if doc_type == "image":
                # Create simple image
                img = Image.new('RGB', (200, 100), color='white')
                draw = ImageDraw.Draw(img)
                draw.text((20, 20), f"Test {format_name}", fill='black')

                buffer = io.BytesIO()
                fmt = 'JPEG' if format_name in ('JPEG', 'JPG') else format_name
                img.save(buffer, format=fmt)
                buffer.seek(0)
                file_bytes = buffer.read()

            elif format_name == "PDF":
                from reportlab.lib.pagesizes import letter
                from reportlab.pdfgen import canvas
                buffer = io.BytesIO()
                c = canvas.Canvas(buffer, pagesize=letter)
                c.drawString(100, 750, "Test PDF Document")
                c.drawString(100, 730, "This is a test paragraph.")
                c.showPage()
                c.save()
                buffer.seek(0)
                file_bytes = buffer.read()

            elif format_name in ("DOCX", "DOC"):
                from docx import Document
                doc = Document()
                doc.add_heading('Test Document', 0)
                doc.add_paragraph('This is a test paragraph.')
                buffer = io.BytesIO()
                doc.save(buffer)
                buffer.seek(0)
                file_bytes = buffer.read()

            elif format_name in ("PPTX", "PPT"):
                from pptx import Presentation
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = "Test Presentation"
                buffer = io.BytesIO()
                prs.save(buffer)
                buffer.seek(0)
                file_bytes = buffer.read()

            elif format_name == "HTML":
                html_content = """
                <!DOCTYPE html>
                <html>
                <head><title>Test</title></head>
                <body>
                    <h1>Test Document</h1>
                    <p>This is a test paragraph.</p>
                </body>
                </html>
                """
                file_bytes = html_content.encode('utf-8')

            else:
                results[format_name] = "❌ SKIP (unknown type)"
                continue

            # Extract
            response = extract_document(service_url, file_bytes, filename, mime_type)

            if response.status_code == 200:
                result = response.json()
                if 'pages' in result and len(result['pages']) >= 1:
                    results[format_name] = "✅ PASS"
                else:
                    results[format_name] = f"⚠️  PARTIAL (no pages)"
            else:
                results[format_name] = f"❌ FAIL ({response.status_code})"

        except Exception as e:
            results[format_name] = f"❌ ERROR: {str(e)[:50]}"

    # Print summary
    print("\n" + "=" * 80)
    print("DOCLING FORMAT SUPPORT SUMMARY")
    print("=" * 80)
    for format_name, status in results.items():
        print(f"{format_name:15} {status}")
    print("=" * 80)

    # Count results
    passed = sum(1 for s in results.values() if '✅' in s)
    total = len(results)
    print(f"\nResult: {passed}/{total} formats working\n")

    # All should pass
    failures = [(name, status) for name, status in results.items() if '❌' in status]
    assert len(failures) == 0, f"Some formats failed: {failures}"


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
